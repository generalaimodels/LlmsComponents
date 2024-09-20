import os
import logging
from typing import Optional
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor
import shutil

import pandas as pd
import json
import yaml
import docx  # python-docx
import pptx  # python-pptx
import fitz  # PyMuPDF

import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords

# Setup NLTK resources quietly
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(name)


def extract_text_from_file(file_path: Path) -> Optional[str]:
    """
    Extracts text content from various file types.
    """
    try:
        if file_path.suffix == '.txt':
            with file_path.open('r', encoding='utf-8') as file:
                return file.read()
        elif file_path.suffix == '.pdf':
            text = ''
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text()
            return text
        elif file_path.suffix == '.docx':
            doc = docx.Document(file_path)
            return '\n'.join([para.text for para in doc.paragraphs])
        elif file_path.suffix == '.pptx':
            prs = pptx.Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return '\n'.join(text_runs)
        elif file_path.suffix == '.json':
            with file_path.open('r', encoding='utf-8') as file:
                data = json.load(file)
            return json.dumps(data)
        elif file_path.suffix in ['.yaml', '.yml']:
            with file_path.open('r', encoding='utf-8') as file:
                data = yaml.safe_load(file)
            return yaml.dump(data)
        elif file_path.suffix == '.csv':
            df = pd.read_csv(file_path)
            return df.to_csv(index=False)
        else:
            logger.warning(f"Unsupported file type: {file_path.suffix}")
            return None
    except Exception as e:
        logger.error(f"Error processing {file_path}: {e}")
        return None


def process_files_in_directory(
    source_dir: Path,
    output_dir: Path,
    num_workers: int = 4
) -> None:
    """
    Recursively processes files in the source directory and saves them as .txt files in the output directory.
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    file_paths = [p for p in source_dir.rglob('*') if p.is_file()]

    def process_file(file_path: Path) -> None:
        relative_path = file_path.relative_to(source_dir)
        output_file_path = output_dir / relative_path.with_suffix('.txt')

        output_file_path.parent.mkdir(parents=True, exist_ok=True)

        text = extract_text_from_file(file_path)
        if text:
            try:
                with output_file_path.open('w', encoding='utf-8') as file:
                    file.write(text)
                logger.info(f"Processed and saved: {output_file_path}")
            except Exception as e:
                logger.error(f"Failed to write {output_file_path}: {e}")

    with ThreadPoolExecutor(max_workers=num_workers) as executor:
        executor.map(process_file, file_paths)


if name == "main":
    import argparse

    parser = argparse.ArgumentParser(description="Convert files to .txt recursively.")
    parser.add_argument("source_dir", type=str, help="Source directory path.")
    parser.add_argument("output_dir", type=str, help="Output directory path.")
    parser.add_argument(
        "--workers",
        type=int,
        default=4,
        help="Number of worker threads for processing files."
    )

    args = parser.parse_args()

    source_directory = Path(args.source_dir)
    output_directory = Path(args.output_dir)

    process_files_in_directory(source_directory, output_directory, args.workers)

import logging
from typing import List

from pathlib import Path
from concurrent.futures import ProcessPoolExecutor

from langchain.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from transformers import AutoTokenizer
from tqdm import tqdm

from langchain.vectorstores import FAISS
from langchain.embeddings.huggingface import HuggingFaceEmbeddings

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(name)

# Constants for embedding
EMBEDDING_MODEL_NAME = "thenlper/gte-base"  # Replace with your preferred model

# Initialize tokenizer
tokenizer = AutoTokenizer.from_pretrained(EMBEDDING_MODEL_NAME)

# Define text splitter
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=512,
    chunk_overlap=64,
    length_function=lambda text: len(tokenizer.encode(text, truncation=True)),
)

def load_and_split_document(file_path: Path) -> List[Document]:
    """
    Loads text from a file and splits it into chunks.
    """
    try:
        with file_path.open('r', encoding='utf-8') as file:
            text = file.read()
        metadata = {"source": str(file_path)}
        texts = text_splitter.split_text(text)
        return [Document(page_content=t, metadata=metadata) for t in texts]
    except Exception as e:
        logger.error(f"Failed to process {file_path}: {e}")
        return []


def build_vector_store(documents: List[Document]) -> FAISS:
    """
    Builds a FAISS vector store from the documents.
    """
    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL_NAME)
    vector_store = FAISS.from_documents(documents, embeddings)
    return vector_store


def search_documents(
    vector_store: FAISS,
    query: str,
    k: int = 5
) -> List[Document]:
    """
    Searches the vector store for the top k documents matching the query.
    """
    results = vector_store.similarity_search(query, k=k)
    return results


if name == "main":
    import argparse

    parser = argparse.ArgumentParser(description="Process text files and retrieve top K matches.")
    parser.add_argument("text_files_dir", type=str, help="Directory containing .txt files.")
    parser.add_argument("query", type=str, help="User query for searching documents.")
    parser.add_argument(
        "--k",
        type=int,
        default=5,
        help="Number of top documents to retrieve."
    )

    args = parser.parse_args()

    text_files_directory = Path(args.text_files_dir)
    all_txt_files = list(text_files_directory.rglob('*.txt'))

    all_documents = []
    with ProcessPoolExecutor() as executor:
        results = executor.map(load_and_split_document, all_txt_files)
        for docs in tqdm(results, total=len(all_txt_files), desc="Processing documents"):
            all_documents.extend(docs)

    logger.info(f"Total chunks created: {len(all_documents)}")

    vector_store = build_vector_store(all_documents)
    search_results = search_documents(vector_store, args.query, k=args.k)

    for idx, doc in enumerate(search_results, 1):
        print(f"\nResult {idx}:")
        print(f"Source: {doc.metadata['source']}")
        print(f"Content:\n{doc.page_content}")


# Instructions to Run:

# - Ensure you have the required packages installed:

#   pip install langchain transformers faiss-cpu tqdm

# - Execute the script from the command line:

#   python script_name.py /path/to/text_files "Your search query here" --k 5


# ---

# Step 3: Save and Load the Vector Store for Reusability

# This code demonstrates how to save the FAISS vector store to disk and load it later for future queries, allowing for efficient reuse without rebuilding the index.

import os
import logging
from pathlib import Path

from langchain.vectorstores import FAISS
from langchain.embeddings.huggingface import HuggingFaceEmbeddings

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(name)

def save_vector_store(vector_store: FAISS, folder_path: str) -> None:
    """
    Saves the FAISS vector store to the specified folder.
    """
    try:
        vector_store.save_local(folder_path)
        logger.info(f"Vector store saved to {folder_path}")
    except Exception as e:
        logger.error(f"Failed to save vector store: {e}")

def load_vector_store(folder_path: str) -> FAISS:
    """
    Loads the FAISS vector store from the specified folder.
    """
    try:
        embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL_NAME)
        vector_store = FAISS.load_local(folder_path, embeddings)
        logger.info(f"Vector store loaded from {folder_path}")
        return vector_store
    except Exception as e:
        logger.error(f"Failed to load vector store: {e}")
        raise e

# # if name == "main":
#     import argparse

#     parser = argparse.ArgumentParser(description="Save and load FAISS vector store.")
#     parser.add_argument("action", choices=["save", "load"], help="Action to perform.")
#     parser.add_argument("folder_path", type=str, help="Path to save to or load from.")
#     args = parser.parse_args()

#     if args.action == "save":
#         # Assume vector_store is already built (e.g., from previous code)
#         save_vector_store(vector_store, args.folder_path)
#     elif args.action == "load":
#         vector_store = load_vector_store(args.folder_path)
