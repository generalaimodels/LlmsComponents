import os
import logging
import shutil
from typing import List, Optional, Dict
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor
import json
import yaml
import pandas as pd
import docx
import pptx
import fitz  # PyMuPDF for PDF processing
import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from datasets import Dataset
from transformers import AutoTokenizer
from tqdm import tqdm
from langchain.docstore.document import Document as LangchainDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores.utils import DistanceStrategy


# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Download necessary NLTK resources quietly
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)

# Constants
EMBEDDING_MODEL_NAME = "thenlper/gte-large"  # Replace with your preferred model
CHUNK_SIZE = 500  # Adjust based on your needs
CHUNK_OVERLAP = 50  # Overlap between chunks for context

def extract_text_from_pdf(file_path: Path) -> str:
    """Extract text from a PDF file."""
    try:
        with fitz.open(file_path) as doc:
            text = ""
            for page in doc:
                text += page.get_text()
        return text
    except Exception as e:
        logger.error(f"Error reading PDF file {file_path}: {e}")
        return ""

def extract_text_from_docx(file_path: Path) -> str:
    """Extract text from a Word document."""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        logger.error(f"Error reading DOCX file {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path: Path) -> str:
    """Extract text from a PowerPoint presentation."""
    try:
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error reading PPTX file {file_path}: {e}")
        return ""

def extract_text_from_json(file_path: Path) -> str:
    """Extract text from a JSON file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return json.dumps(data)
    except Exception as e:
        logger.error(f"Error reading JSON file {file_path}: {e}")
        return ""

def extract_text_from_yaml(file_path: Path) -> str:
    """Extract text from a YAML file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = yaml.safe_load(f)
        return yaml.dump(data)
    except Exception as e:
        logger.error(f"Error reading YAML file {file_path}: {e}")
        return ""

def extract_text_from_csv(file_path: Path) -> str:
    """Extract text from a CSV file."""
    try:
        df = pd.read_csv(file_path)
        return df.to_csv(index=False)
    except Exception as e:
        logger.error(f"Error reading CSV file {file_path}: {e}")
        return ""

def extract_text(file_path: Path) -> str:
    """Extract text from a file based on its extension."""
    ext = file_path.suffix.lower()
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.json':
        return extract_text_from_json(file_path)
    elif ext in ['.yml', '.yaml']:
        return extract_text_from_yaml(file_path)
    elif ext == '.csv':
        return extract_text_from_csv(file_path)
    elif ext == '.txt':
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading TXT file {file_path}: {e}")
            return ""
    else:
        logger.warning(f"Unsupported file type: {file_path}")
        return ""

def process_files(input_dir: Path, output_dir: Path):
    """Recursively process files in input_dir and save text files to output_dir."""
    if not output_dir.exists():
        output_dir.mkdir(parents=True)

    def process_file(file_path: Path):
        text = extract_text(file_path)
        if text:
            relative_path = file_path.relative_to(input_dir)
            output_file = output_dir.joinpath(relative_path).with_suffix('.txt')
            output_file.parent.mkdir(parents=True, exist_ok=True)
            try:
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write(text)
                logger.info(f"Processed and saved text from {file_path} to {output_file}")
            except Exception as e:
                logger.error(f"Error writing to file {output_file}: {e}")

    files = [p for p in input_dir.rglob('') if p.is_file()]
    with ThreadPoolExecutor() as executor:
        executor.map(process_file, files)


def load_documents(txt_folder: Path) -> List[LangchainDocument]:
    """Load text files and return as LangchainDocument objects."""
    documents = []
    files = list(txt_folder.rglob('*.txt'))  # Corrected pattern here
    for file_path in tqdm(files, desc="Loading documents"):
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            metadata = {
                'source': str(file_path),
                'file_name': file_path.name,
                'file_path': str(file_path)
            }
            document = LangchainDocument(page_content=text, metadata=metadata)
            documents.append(document)
        except Exception as e:
            logger.error(f"Error loading document {file_path}: {e}")
    return documents


def split_documents(documents: List[LangchainDocument]) -> List[LangchainDocument]:
    """Split documents into smaller chunks."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", " ", ""]
    )
    split_docs = []
    for doc in tqdm(documents, desc="Splitting documents"):
        splits = text_splitter.split_text(doc.page_content)
        for i, content in enumerate(splits):
            metadata = doc.metadata.copy()
            metadata['chunk'] = i
            split_docs.append(LangchainDocument(page_content=content, metadata=metadata))
    return split_docs

def create_embeddings(documents: List[LangchainDocument]) -> FAISS:
    """Create embeddings and build FAISS index."""
    embeddings = HuggingFaceEmbeddings(
        model_name=EMBEDDING_MODEL_NAME,
        model_kwargs={'device': 'cpu'}  # Change to 'cuda' if using GPU
    )
    vector_store = FAISS.from_documents(
        documents,
        embedding=embeddings,
        distance_strategy=DistanceStrategy.EUCLIDEAN_DISTANCE
    )
    return vector_store

def save_vector_store(vector_store: FAISS, index_path: Path, metadata_path: Path):
    """Save the FAISS index and metadata."""
    vector_store.save_local(str(index_path), str(metadata_path))
    logger.info(f"Vector store saved to {index_path} and {metadata_path}")

def load_vector_store(index_path: Path, metadata_path: Path) -> FAISS:
    """Load the FAISS index and metadata."""
    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL_NAME)
    vector_store = FAISS.load_local(str(index_path), embeddings, str(metadata_path))
    logger.info(f"Vector store loaded from {index_path} and {metadata_path}")
    return vector_store

def query_vector_store(vector_store: FAISS, query: str, k: int = 5) -> List[Dict]:
    """Query the vector store and get top k documents."""
    docs = vector_store.similarity_search(query, k=k)
    results = []
    for doc in docs:
        result = {
            'content': doc.page_content,
            'metadata': doc.metadata
        }
        results.append(result)
    return results



def main():
    # Define input and output directories
    # input_dir = Path(r'E:\LLMS\Fine-tuning\test')  # Replace with your input directory
    txt_output_dir = Path(r"E:\LLMS\Fine-tuning\output1")  # Replace with your output directory

    # Step 1: Process files and convert to .txt
    # process_files(input_dir, txt_output_dir)

    # Step 2: Load documents and build vector database
    documents = load_documents(txt_output_dir)
    if not documents:
        logger.error("No documents found to process. Please check the input directory and file patterns.")
        return

    split_docs = split_documents(documents)
    if not split_docs:
        logger.error("No split documents created. Please check the text splitter configuration.")
        return

    logger.info(f"Loaded {len(documents)} documents and split into {len(split_docs)} chunks")

    vector_store = create_embeddings(split_docs)

    # Save the vector store
    index_path = Path('faiss_index')
    metadata_path = Path('faiss_index_metadata')
    save_vector_store(vector_store, index_path, metadata_path)

    # Load the vector store (optional, to demonstrate loading)
    # vector_store = load_vector_store(index_path, metadata_path)

    # Query the vector store
    user_query = "Tweakable, contributor-friendly over abstraction"  # Replace with the user's query
    top_k = 5  # Replace with desired number of results
    results = query_vector_store(vector_store, user_query, k=top_k)

    # Display results
    for i, result in enumerate(results):
        print(f"Result {i + 1}:")
        print(f"Content:\n{result['content']}\n")
        print(f"Metadata:\n{result['metadata']}\n")
        print("=" * 50)





if __name__ == "__main__":
   main()